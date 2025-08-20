#!/usr/bin/env python3
"""
PowerPoint to Markdown Converter

This script extracts text content with layout information from PPTX files
in a specified folder, orders them, and combines them into a single markdown file.

Usage:
    python extract_pptx_to_markdown.py <input_folder> [output_file]

Requirements:
    pip install python-pptx
"""

import os
import sys
import re
import argparse
from pathlib import Path
from typing import List, Dict, Tuple

try:
    from pptx import Presentation
    from pptx.shapes.base import BaseShape
    from pptx.shapes.autoshape import Shape
    from pptx.shapes.picture import Picture
    from pptx.table import Table
    from pptx.shapes.group import GroupShape
except ImportError:
    print("Error: python-pptx library not found.")
    print("Please install it using: pip install python-pptx")
    sys.exit(1)


class PPTXExtractor:
    """Extract text and layout information from PowerPoint files"""
    
    def __init__(self):
        self.slide_counter = 0
        
    def extract_text_from_shape(self, shape: BaseShape, ignore_images: bool = False) -> str:
        """Extract text from a shape, handling different shape types"""
        text_content = []
        
        try:
            if hasattr(shape, 'text') and shape.text.strip():
                # Regular text shape
                text_content.append(shape.text.strip())
                
            elif hasattr(shape, 'table') and shape.table is not None:
                # Table shape
                table_content = self.extract_table_content(shape.table)
                if table_content:
                    text_content.append(table_content)
                    
            elif isinstance(shape, GroupShape):
                # Group shape - recursively extract from grouped shapes
                for grouped_shape in shape.shapes:
                    grouped_text = self.extract_text_from_shape(grouped_shape, ignore_images)
                    if grouped_text:
                        text_content.append(grouped_text)
                        
            elif isinstance(shape, Picture) and not ignore_images:
                # Picture shape - add placeholder only if not ignoring images
                text_content.append("[Image]")
                
        except Exception as e:
            # Skip problematic shapes and continue processing
            pass
            
        return "\n".join(text_content)
    
    def extract_table_content(self, table: Table) -> str:
        """Extract content from a table and format as markdown table"""
        if not table.rows:
            return ""
            
        table_lines = []
        
        # Extract header row
        if table.rows:
            header_cells = []
            for cell in table.rows[0].cells:
                header_cells.append(cell.text.strip() or " ")
            table_lines.append("| " + " | ".join(header_cells) + " |")
            table_lines.append("|" + "---|" * len(header_cells))
            
            # Extract data rows
            for row in table.rows[1:]:
                data_cells = []
                for cell in row.cells:
                    data_cells.append(cell.text.strip() or " ")
                table_lines.append("| " + " | ".join(data_cells) + " |")
                
        return "\n".join(table_lines)
    
    def extract_slide_content(self, slide, ignore_images: bool = False) -> Dict[str, str]:
        """Extract content from a single slide with layout information"""
        self.slide_counter += 1
        
        slide_content = {
            'slide_number': self.slide_counter,
            'title': '',
            'content': [],
            'notes': ''
        }
        
        # Extract slide title (usually the first text box or title placeholder)
        title_found = False
        first_shape_top = None
        
        # Get the top position of the first shape for comparison
        if slide.shapes:
            first_shape_top = slide.shapes[0].top
        
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'text') and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Try to identify title based on position and formatting
                    is_likely_title = (
                        not title_found and 
                        (first_shape_top is None or shape.top <= first_shape_top + 100000) and
                        len(text.split('\n')) == 1 and len(text) < 100
                    )
                    
                    if is_likely_title:
                        slide_content['title'] = text
                        title_found = True
                    else:
                        # Regular content
                        extracted_text = self.extract_text_from_shape(shape, ignore_images)
                        if extracted_text and extracted_text != slide_content['title']:
                            slide_content['content'].append(extracted_text)
                else:
                    # Non-text shapes (tables, images, etc.)
                    extracted_text = self.extract_text_from_shape(shape, ignore_images)
                    if extracted_text:
                        slide_content['content'].append(extracted_text)
            except Exception as e:
                # Skip problematic shapes and continue processing
                continue
        
        # Extract slide notes if available
        if hasattr(slide, 'notes_slide') and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_content['notes'] = notes_text
                
        return slide_content
    
    def extract_presentation(self, pptx_path: str, ignore_images: bool = False) -> Dict[str, any]:
        """Extract all content from a PowerPoint presentation"""
        try:
            prs = Presentation(pptx_path)
            
            presentation_data = {
                'filename': Path(pptx_path).name,
                'title': Path(pptx_path).stem,
                'slides': [],
                'slide_count': len(prs.slides)
            }
            
            # Reset slide counter for each presentation
            self.slide_counter = 0
            
            for slide in prs.slides:
                slide_content = self.extract_slide_content(slide, ignore_images)
                presentation_data['slides'].append(slide_content)
                
            return presentation_data
            
        except Exception as e:
            print(f"Error processing {pptx_path}: {str(e)}")
            return None


def find_pptx_files(folder_path: str) -> List[str]:
    """Find all PPTX files in the specified folder"""
    folder = Path(folder_path)
    if not folder.exists() or not folder.is_dir():
        print(f"Error: Folder '{folder_path}' does not exist or is not a directory")
        return []
    
    pptx_files = list(folder.glob("*.pptx"))
    pptx_files.extend(list(folder.glob("*.PPTX")))
    
    # Sort files naturally (handling numbers in filenames)
    def natural_sort_key(filename):
        return [int(text) if text.isdigit() else text.lower() for text in re.split(r'(\d+)', str(filename))]
    
    pptx_files.sort(key=natural_sort_key)
    return [str(f) for f in pptx_files]


def format_slide_as_markdown(slide_data: Dict[str, any], presentation_title: str) -> str:
    """Format a single slide as markdown"""
    markdown_lines = []
    
    # Slide header
    if slide_data['title']:
        markdown_lines.append(f"### {slide_data['title']}")
    else:
        markdown_lines.append(f"### Slide {slide_data['slide_number']}")
    
    markdown_lines.append("")
    
    # Slide content
    for content_item in slide_data['content']:
        if content_item.strip():
            # Check if this looks like code content (extremely restrictive - only executable Python syntax)
            is_code_block = (
                # Actual variable assignments with specific patterns
                (re.search(r'\w+\s*=\s*np\.[a-zA-Z_]\w*\(', content_item)) or
                # Print statements with parentheses and quotes
                (re.search(r'print\s*\(["\'].*["\']\)', content_item)) or
                # Import statements with specific modules
                (re.search(r'^\s*import\s+(numpy|torch|matplotlib)', content_item, re.MULTILINE)) or
                (re.search(r'^\s*from\s+(numpy|torch|matplotlib)', content_item, re.MULTILINE)) or
                # Function definitions with return statements or function calls
                (re.search(r'def\s+\w+\s*\([^)]*\)\s*:', content_item) and ('return ' in content_item or 'print(' in content_item)) or
                # Mathematical expressions with actual operators and parentheses
                (re.search(r'\w+\s*=\s*\w+\s*[+\-*/]\s*\w+', content_item) and '(' in content_item and ')' in content_item)
            )
            
            if is_code_block:
                # Format as code block
                markdown_lines.append("```python")
                markdown_lines.append(content_item)
                markdown_lines.append("```")
                markdown_lines.append("")
            else:
                # Handle bullet points and regular formatting
                lines = content_item.split('\n')
                for line in lines:
                    line = line.strip()
                    if line:
                        # Escape lines that start with # to prevent markdown header interpretation
                        if line.startswith('#') and not line.startswith('##'):
                            line = '`' + line + '`'  # Format as inline code
                        
                        # Convert bullet points
                        if line.startswith('•') or line.startswith('-'):
                            markdown_lines.append(f"- {line[1:].strip()}")
                        else:
                            markdown_lines.append(line)
                markdown_lines.append("")
    
    # Add notes if available
    if slide_data['notes']:
        markdown_lines.append("**Speaker Notes:**")
        markdown_lines.append(slide_data['notes'])
        markdown_lines.append("")
    
    return "\n".join(markdown_lines)


def create_combined_markdown(presentations: List[Dict[str, any]], output_file: str):
    """Create a combined markdown file from all presentations"""
    markdown_content = []
    
    # Document header
    markdown_content.append("# Combined Presentation Content")
    markdown_content.append("")
    markdown_content.append(f"*Generated from {len(presentations)} PowerPoint presentations*")
    markdown_content.append("")
    
    # Table of contents
    markdown_content.append("## Table of Contents")
    markdown_content.append("")
    for i, pres in enumerate(presentations, 1):
        markdown_content.append(f"{i}. [{pres['title']}](#{pres['title'].lower().replace(' ', '-')}) ({pres['slide_count']} slides)")
    markdown_content.append("")
    markdown_content.append("---")
    markdown_content.append("")
    
    # Process each presentation
    for pres in presentations:
        # Presentation header
        markdown_content.append(f"## {pres['title']}")
        markdown_content.append("")
        markdown_content.append(f"*Source: {pres['filename']}*")
        markdown_content.append("")
        
        # Process each slide
        for slide in pres['slides']:
            slide_markdown = format_slide_as_markdown(slide, pres['title'])
            markdown_content.append(slide_markdown)
        
        markdown_content.append("---")
        markdown_content.append("")
    
    # Write to file
    try:
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write("\n".join(markdown_content))
        print(f"✅ Combined markdown file created: {output_file}")
    except Exception as e:
        print(f"❌ Error writing to {output_file}: {str(e)}")


def main():
    """Main function"""
    parser = argparse.ArgumentParser(
        description="Extract text and layout information from PPTX files and combine into a single Markdown file",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""Examples:
  python extract_pptx_to_markdown.py /path/to/pptx/folder
  python extract_pptx_to_markdown.py /path/to/pptx/folder output.md
  python extract_pptx_to_markdown.py /path/to/pptx/folder output.md --ignore-images
  python extract_pptx_to_markdown.py --help"""
    )
    
    parser.add_argument(
        'input_folder',
        help='Path to folder containing PPTX files'
    )
    
    parser.add_argument(
        'output_file',
        nargs='?',
        default='combined_presentations.md',
        help='Output markdown file (default: combined_presentations.md)'
    )
    
    parser.add_argument(
        '--ignore-images',
        action='store_true',
        help='Ignore image placeholders in the output'
    )
    
    args = parser.parse_args()
    
    input_folder = args.input_folder
    output_file = args.output_file
    ignore_images = args.ignore_images
    
    print(f"🔍 Searching for PPTX files in: {input_folder}")
    
    # Find all PPTX files
    pptx_files = find_pptx_files(input_folder)
    
    if not pptx_files:
        print("❌ No PPTX files found in the specified folder")
        sys.exit(1)
    
    print(f"📁 Found {len(pptx_files)} PPTX files:")
    for i, file in enumerate(pptx_files, 1):
        print(f"  {i}. {Path(file).name}")
    
    if ignore_images:
        print("🖼️  Image placeholders will be ignored")
    
    print("")
    
    # Extract content from all presentations
    extractor = PPTXExtractor()
    presentations = []
    
    for pptx_file in pptx_files:
        print(f"📖 Processing: {Path(pptx_file).name}")
        presentation_data = extractor.extract_presentation(pptx_file, ignore_images)
        
        if presentation_data:
            presentations.append(presentation_data)
            print(f"   ✅ Extracted {presentation_data['slide_count']} slides")
        else:
            print(f"   ❌ Failed to process file")
    
    if not presentations:
        print("❌ No presentations were successfully processed")
        sys.exit(1)
    
    print("")
    print(f"📝 Creating combined markdown file: {output_file}")
    
    # Create combined markdown
    create_combined_markdown(presentations, output_file)
    
    # Summary
    total_slides = sum(pres['slide_count'] for pres in presentations)
    print("")
    print("📊 Summary:")
    print(f"   - Presentations processed: {len(presentations)}")
    print(f"   - Total slides extracted: {total_slides}")
    print(f"   - Output file: {output_file}")
    if ignore_images:
        print(f"   - Images: Ignored")
    print("")
    print("🎉 Conversion completed successfully!")


if __name__ == "__main__":
    main()