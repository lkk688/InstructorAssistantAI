#!/usr/bin/env python3
"""
PowerPoint Structure Analyzer and Pandoc Config Generator

This script analyzes existing PPTX files to understand their organization,
templates, and styling patterns, then generates optimized pandoc_revealjs_config.yaml
files based on the discovered structure.

Usage:
    python analyze_pptx_structure.py <input_folder> [config_output]

Requirements:
    pip install python-pptx pyyaml
"""

import os
import sys
import yaml
import re
from pathlib import Path
from typing import List, Dict, Tuple, Optional
from collections import Counter, defaultdict

try:
    from pptx import Presentation
    from pptx.shapes.base import BaseShape
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx library not found.")
    print("Please install it using: pip install python-pptx")
    sys.exit(1)

try:
    import yaml
except ImportError:
    print("Error: PyYAML library not found.")
    print("Please install it using: pip install pyyaml")
    sys.exit(1)


class PPTXStructureAnalyzer:
    """Analyze PowerPoint structure and generate pandoc configurations"""
    
    def __init__(self):
        self.presentations = []
        self.structure_stats = {
            'slide_counts': [],
            'title_patterns': [],
            'content_types': defaultdict(int),
            'color_schemes': [],
            'font_families': [],
            'slide_layouts': defaultdict(int),
            'header_levels': defaultdict(int),
            'has_speaker_notes': 0,
            'has_tables': 0,
            'has_images': 0,
            'has_charts': 0
        }
    
    def analyze_color(self, color_obj) -> Optional[str]:
        """Extract color information from PowerPoint color objects"""
        try:
            if hasattr(color_obj, 'rgb'):
                rgb = color_obj.rgb
                return f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
        except:
            pass
        return None
    
    def analyze_font(self, font_obj) -> Dict[str, any]:
        """Extract font information"""
        font_info = {}
        try:
            if hasattr(font_obj, 'name') and font_obj.name:
                font_info['family'] = font_obj.name
            if hasattr(font_obj, 'size') and font_obj.size:
                font_info['size'] = font_obj.size.pt
            if hasattr(font_obj, 'color'):
                color = self.analyze_color(font_obj.color)
                if color:
                    font_info['color'] = color
        except:
            pass
        return font_info
    
    def analyze_slide_layout(self, slide) -> str:
        """Determine slide layout type based on content structure"""
        shapes = list(slide.shapes)
        text_shapes = [s for s in shapes if hasattr(s, 'text') and s.text.strip()]
        
        if len(text_shapes) == 0:
            return "blank"
        elif len(text_shapes) == 1:
            return "title_only"
        elif len(text_shapes) == 2:
            # Check if first is title-like
            first_text = text_shapes[0].text.strip()
            if len(first_text.split('\n')) == 1 and len(first_text) < 100:
                return "title_content"
        elif len(text_shapes) > 2:
            return "multi_content"
        
        # Check for specific content types
        has_table = any(hasattr(s, 'table') for s in shapes)
        has_image = any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in shapes)
        has_chart = any(s.shape_type == MSO_SHAPE_TYPE.CHART for s in shapes)
        
        if has_table:
            return "table_layout"
        elif has_image:
            return "image_layout"
        elif has_chart:
            return "chart_layout"
        
        return "standard"
    
    def analyze_content_structure(self, text: str) -> Dict[str, int]:
        """Analyze text content structure"""
        structure = {
            'paragraphs': len([p for p in text.split('\n\n') if p.strip()]),
            'bullet_points': len(re.findall(r'^[•\-\*]', text, re.MULTILINE)),
            'numbered_items': len(re.findall(r'^\d+\.', text, re.MULTILINE)),
            'questions': len(re.findall(r'\?', text)),
            'exclamations': len(re.findall(r'!', text)),
            'code_blocks': len(re.findall(r'```|`[^`]+`', text)),
            'urls': len(re.findall(r'https?://\S+', text))
        }
        return structure
    
    def analyze_presentation(self, pptx_path: str) -> Dict[str, any]:
        """Analyze a single PowerPoint presentation"""
        try:
            prs = Presentation(pptx_path)
            
            analysis = {
                'filename': Path(pptx_path).name,
                'slide_count': len(prs.slides),
                'slides': [],
                'themes': [],
                'fonts': [],
                'colors': [],
                'has_master_slides': len(prs.slide_masters) > 0
            }
            
            # Analyze each slide
            for i, slide in enumerate(prs.slides):
                slide_analysis = {
                    'slide_number': i + 1,
                    'layout_type': self.analyze_slide_layout(slide),
                    'text_content': [],
                    'has_notes': False,
                    'content_structure': {},
                    'fonts_used': [],
                    'colors_used': []
                }
                
                # Analyze slide content
                all_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        slide_analysis['text_content'].append(text)
                        all_text += text + "\n"
                        
                        # Analyze text formatting
                        if hasattr(shape, 'text_frame'):
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.font:
                                        font_info = self.analyze_font(run.font)
                                        if font_info:
                                            slide_analysis['fonts_used'].append(font_info)
                
                # Analyze content structure
                if all_text:
                    slide_analysis['content_structure'] = self.analyze_content_structure(all_text)
                
                # Check for speaker notes
                if hasattr(slide, 'notes_slide') and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_analysis['has_notes'] = True
                        self.structure_stats['has_speaker_notes'] += 1
                
                # Check for tables, images, charts
                for shape in slide.shapes:
                    if hasattr(shape, 'table'):
                        self.structure_stats['has_tables'] += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        self.structure_stats['has_images'] += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        self.structure_stats['has_charts'] += 1
                
                analysis['slides'].append(slide_analysis)
                self.structure_stats['slide_layouts'][slide_analysis['layout_type']] += 1
            
            # Update global statistics
            self.structure_stats['slide_counts'].append(analysis['slide_count'])
            
            return analysis
            
        except Exception as e:
            print(f"Error analyzing {pptx_path}: {str(e)}")
            return None
    
    def generate_pandoc_config(self, output_path: str):
        """Generate optimized pandoc configuration based on analysis"""
        if not self.presentations:
            print("No presentations analyzed. Cannot generate configuration.")
            return
        
        # Calculate statistics
        total_slides = sum(self.structure_stats['slide_counts'])
        avg_slides = total_slides / len(self.structure_stats['slide_counts'])
        most_common_layout = max(self.structure_stats['slide_layouts'].items(), key=lambda x: x[1])[0]
        
        # Determine optimal settings based on analysis
        config = {
            '# Generated Configuration': 'Based on PPTX structure analysis',
            'title': '"Presentation from PPTX Analysis"',
            'author': '"Generated Author"',
            'date': '"2024"',
            
            # reveal.js settings optimized for discovered content
            'revealjs-url': '"https://unpkg.com/reveal.js@4.3.1/"',
            'theme': self._determine_optimal_theme(),
            'transition': self._determine_optimal_transition(),
            'backgroundTransition': '"fade"',
            'hash': True,
            'controls': True,
            'progress': True,
            'center': self._should_center_content(),
            'touch': True,
            'loop': False,
            'rtl': False,
            'navigationMode': '"default"',
            'previewLinks': False,
            'hideAddressBar': True,
            
            # Display settings based on content density
            'width': 1280,
            'height': 720,
            'margin': self._calculate_optimal_margin(),
            'minScale': 0.2,
            'maxScale': 1.5,
            
            # Slide structure
            'slide-level': self._determine_slide_level(),
            'incremental': self._should_use_incremental(),
            
            # Content processing
            'highlight-style': '"github"',
            'mathjax': self._needs_math_support(),
            'speaker-notes': self.structure_stats['has_speaker_notes'] > 0,
            
            # Custom CSS based on discovered patterns
            'css': [self._generate_custom_css()]
        }
        
        # Add variables
        config['variables'] = {
            'fontsize': '"18pt"',
            'mainfont': self._determine_main_font(),
            'monofont': '"Consolas"'
        }
        
        # Write configuration
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                # Write header comment
                f.write(f"# Pandoc reveal.js Configuration\n")
                f.write(f"# Generated from analysis of {len(self.presentations)} PPTX files\n")
                f.write(f"# Total slides analyzed: {total_slides}\n")
                f.write(f"# Average slides per presentation: {avg_slides:.1f}\n")
                f.write(f"# Most common layout: {most_common_layout}\n")
                f.write(f"# Speaker notes found: {self.structure_stats['has_speaker_notes'] > 0}\n")
                f.write(f"# Tables found: {self.structure_stats['has_tables'] > 0}\n")
                f.write(f"# Images found: {self.structure_stats['has_images'] > 0}\n")
                f.write(f"# Charts found: {self.structure_stats['has_charts'] > 0}\n\n")
                
                # Write YAML configuration
                yaml.dump(config, f, default_flow_style=False, allow_unicode=True, sort_keys=False)
                
            print(f"✅ Generated pandoc configuration: {output_path}")
            
        except Exception as e:
            print(f"❌ Error writing configuration: {str(e)}")
    
    def _determine_optimal_theme(self) -> str:
        """Determine the best theme based on content analysis"""
        # Default themes based on content type
        if self.structure_stats['has_charts'] > 0:
            return '"white"'  # Better for charts and data
        elif self.structure_stats['has_images'] > 0:
            return '"black"'  # Better contrast for images
        else:
            return '"league"'  # Good general purpose theme
    
    def _determine_optimal_transition(self) -> str:
        """Determine transition based on presentation style"""
        avg_slides = sum(self.structure_stats['slide_counts']) / len(self.structure_stats['slide_counts'])
        if avg_slides > 50:
            return '"fade"'  # Faster for long presentations
        else:
            return '"slide"'  # More engaging for shorter presentations
    
    def _should_center_content(self) -> bool:
        """Determine if content should be centered"""
        # Center if mostly title slides or short content
        title_layouts = self.structure_stats['slide_layouts'].get('title_only', 0)
        total_layouts = sum(self.structure_stats['slide_layouts'].values())
        return title_layouts / total_layouts > 0.3 if total_layouts > 0 else True
    
    def _calculate_optimal_margin(self) -> float:
        """Calculate optimal margin based on content density"""
        if self.structure_stats['has_tables'] > 0:
            return 0.05  # Less margin for tables
        elif self.structure_stats['slide_layouts'].get('multi_content', 0) > 0:
            return 0.08  # Moderate margin for multi-content
        else:
            return 0.1   # Standard margin
    
    def _determine_slide_level(self) -> int:
        """Determine appropriate slide level based on header analysis"""
        # This would need more sophisticated header analysis
        # For now, default to level 3 (###) which matches PPTX extraction
        return 3
    
    def _should_use_incremental(self) -> bool:
        """Determine if incremental reveals should be used"""
        bullet_heavy = sum(1 for slide in self.presentations for s in slide.get('slides', []) 
                          if s.get('content_structure', {}).get('bullet_points', 0) > 3)
        return bullet_heavy > len(self.presentations) * 0.3
    
    def _needs_math_support(self) -> bool:
        """Check if mathematical content is present"""
        # Simple heuristic - look for math-like patterns
        math_indicators = ['∑', '∫', '∂', '≤', '≥', '±', '×', '÷', 'equation', 'formula']
        for pres in self.presentations:
            for slide in pres.get('slides', []):
                for text in slide.get('text_content', []):
                    if any(indicator in text.lower() for indicator in math_indicators):
                        return True
        return False
    
    def _determine_main_font(self) -> str:
        """Determine the most commonly used font"""
        font_counter = Counter()
        for pres in self.presentations:
            for slide in pres.get('slides', []):
                for font_info in slide.get('fonts_used', []):
                    if 'family' in font_info:
                        font_counter[font_info['family']] += 1
        
        if font_counter:
            most_common = font_counter.most_common(1)[0][0]
            return f'"{most_common}"'
        return '"Arial"'
    
    def _generate_custom_css(self) -> str:
        """Generate custom CSS based on discovered patterns"""
        css_rules = []
        
        # Base styling
        css_rules.append("""
        /* Generated CSS based on PPTX analysis */
        .reveal {
            font-family: "Arial", sans-serif;
        }
        """)
        
        # Table styling if tables are present
        if self.structure_stats['has_tables'] > 0:
            css_rules.append("""
        .reveal table {
            font-size: 0.7em;
            border-collapse: collapse;
            margin: 1em auto;
        }
        .reveal table th {
            background-color: #2E86AB;
            color: white;
            padding: 0.5em;
        }
        .reveal table td {
            padding: 0.4em;
            border: 1px solid #ddd;
        }
        """)
        
        # Image styling if images are present
        if self.structure_stats['has_images'] > 0:
            css_rules.append("""
        .reveal img {
            max-width: 80%;
            max-height: 60vh;
            margin: 0.5em auto;
            border-radius: 8px;
        }
        """)
        
        # Speaker notes styling if notes are present
        if self.structure_stats['has_speaker_notes'] > 0:
            css_rules.append("""
        .reveal .speaker-notes {
            background: rgba(162, 59, 114, 0.2);
            border-left: 4px solid #A23B72;
            padding: 1em;
            margin: 1em 0;
            font-size: 0.8em;
        }
        """)
        
        return "".join(css_rules)


def find_pptx_files(folder_path: str) -> List[str]:
    """Find all PPTX files in the specified folder"""
    folder = Path(folder_path)
    if not folder.exists() or not folder.is_dir():
        print(f"Error: Folder '{folder_path}' does not exist or is not a directory")
        return []
    
    pptx_files = list(folder.glob("*.pptx"))
    pptx_files.extend(list(folder.glob("*.PPTX")))
    
    # Filter out temporary files
    pptx_files = [f for f in pptx_files if not f.name.startswith('~$')]
    
    return [str(f) for f in sorted(pptx_files)]


def main():
    """Main function"""
    if len(sys.argv) < 2:
        print("Usage: python analyze_pptx_structure.py <input_folder> [config_output]")
        print("")
        print("Arguments:")
        print("  input_folder  : Path to folder containing PPTX files")
        print("  config_output : Output YAML config file (default: pandoc_config_generated.yaml)")
        sys.exit(1)
    
    input_folder = sys.argv[1]
    config_output = sys.argv[2] if len(sys.argv) > 2 else "pandoc_config_generated.yaml"
    
    print(f"🔍 Analyzing PPTX files in: {input_folder}")
    
    # Find all PPTX files
    pptx_files = find_pptx_files(input_folder)
    
    if not pptx_files:
        print("❌ No PPTX files found in the specified folder")
        sys.exit(1)
    
    print(f"📁 Found {len(pptx_files)} PPTX files:")
    for i, file in enumerate(pptx_files, 1):
        print(f"  {i}. {Path(file).name}")
    print("")
    
    # Analyze presentations
    analyzer = PPTXStructureAnalyzer()
    
    for pptx_file in pptx_files:
        print(f"📖 Analyzing: {Path(pptx_file).name}")
        analysis = analyzer.analyze_presentation(pptx_file)
        
        if analysis:
            analyzer.presentations.append(analysis)
            print(f"   ✅ Analyzed {analysis['slide_count']} slides")
        else:
            print(f"   ❌ Failed to analyze file")
    
    if not analyzer.presentations:
        print("❌ No presentations were successfully analyzed")
        sys.exit(1)
    
    print("")
    print(f"📊 Analysis Summary:")
    total_slides = sum(analyzer.structure_stats['slide_counts'])
    print(f"   - Presentations: {len(analyzer.presentations)}")
    print(f"   - Total slides: {total_slides}")
    print(f"   - Avg slides per presentation: {total_slides/len(analyzer.presentations):.1f}")
    print(f"   - Speaker notes found: {analyzer.structure_stats['has_speaker_notes']}")
    print(f"   - Tables found: {analyzer.structure_stats['has_tables']}")
    print(f"   - Images found: {analyzer.structure_stats['has_images']}")
    print(f"   - Charts found: {analyzer.structure_stats['has_charts']}")
    
    if analyzer.structure_stats['slide_layouts']:
        print(f"   - Most common layout: {max(analyzer.structure_stats['slide_layouts'].items(), key=lambda x: x[1])[0]}")
    
    print("")
    print(f"⚙️  Generating pandoc configuration: {config_output}")
    
    # Generate configuration
    analyzer.generate_pandoc_config(config_output)
    
    print("")
    print("🎉 Analysis and configuration generation completed!")
    print("")
    print("💡 Next steps:")
    print(f"   1. Review the generated config: {config_output}")
    print(f"   2. Extract PPTX content: python extract_pptx_to_markdown.py {input_folder}")
    print(f"   3. Generate presentation: pandoc -d {config_output} combined_presentations.md -o presentation.html")


if __name__ == "__main__":
    main()