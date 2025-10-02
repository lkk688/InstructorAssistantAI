#!/usr/bin/env python3
"""
Enhanced PowerPoint to Markdown Converter with Layout Preservation and OCR

This script extracts text content with complete layout information, images, and OCR
capabilities from PPTX files, preserving spatial relationships and converting
equations to LaTeX format.

Features:
- Layout preservation (position, size, formatting)
- Image extraction and organization
- OCR for text in images
- Equation detection and LaTeX conversion
- Presentation-ready markdown output

Usage:
    python enhanced_pptx_extractor.py <input_folder> [output_file] [--options]

Requirements:
    pip install python-pptx pillow easyocr pix2tex opencv-python numpy
"""

import os
import sys
import re
import json
import argparse
import logging
from pathlib import Path
from typing import List, Dict, Tuple, Optional, Any
from dataclasses import dataclass, asdict
import hashlib
import base64
from io import BytesIO

# Progress bar support
try:
    from tqdm import tqdm
    TQDM_AVAILABLE = True
except ImportError:
    TQDM_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.shapes.base import BaseShape
    from pptx.shapes.autoshape import Shape
    from pptx.shapes.picture import Picture
    from pptx.table import Table
    from pptx.shapes.group import GroupShape
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches, Pt
except ImportError:
    print("Error: python-pptx library not found.")
    print("Please install it using: pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
    import numpy as np
    import cv2
except ImportError:
    print("Error: PIL, numpy, or opencv-python not found.")
    print("Please install them using: pip install pillow numpy opencv-python")
    sys.exit(1)

# Optional OCR imports - will be handled gracefully if not available
try:
    import easyocr
    EASYOCR_AVAILABLE = True
except ImportError:
    EASYOCR_AVAILABLE = False

try:
    from pix2tex.cli import LatexOCR
    PIX2TEX_AVAILABLE = True
except ImportError:
    PIX2TEX_AVAILABLE = False


@dataclass
class LayoutInfo:
    """Store layout information for shapes"""
    left: float
    top: float
    width: float
    height: float
    rotation: float = 0.0
    z_order: int = 0
    
    def to_css_style(self) -> str:
        """Convert layout info to CSS-like positioning"""
        return f"position: absolute; left: {self.left:.1f}px; top: {self.top:.1f}px; width: {self.width:.1f}px; height: {self.height:.1f}px;"


@dataclass
class TextFormatting:
    """Store text formatting information"""
    font_name: str = ""
    font_size: float = 12.0
    bold: bool = False
    italic: bool = False
    underline: bool = False
    color: str = "#000000"
    alignment: str = "left"


@dataclass
class ShapeData:
    """Complete shape data with layout and content"""
    shape_id: str
    shape_type: str
    layout: LayoutInfo
    content: str = ""
    formatting: Optional[TextFormatting] = None
    image_path: str = ""
    ocr_text: str = ""
    latex_equations: List[str] = None
    
    def __post_init__(self):
        if self.latex_equations is None:
            self.latex_equations = []


class EnhancedPPTXExtractor:
    """Enhanced extractor with layout preservation and OCR capabilities"""
    
    def __init__(self, output_dir: str = "extracted_presentations"):
        self.output_dir = Path(output_dir)
        self.images_dir = self.output_dir / "images"
        self.slide_counter = 0
        self.presentation_counter = 0
        
        # Create output directories
        self.output_dir.mkdir(exist_ok=True)
        self.images_dir.mkdir(exist_ok=True)
        
        # Setup logging first
        self._setup_logging()
        
        # Initialize OCR models
        self.ocr_reader = None
        self.latex_ocr = None
        self._init_ocr_models()
    
    def _setup_logging(self):
        """Setup logging configuration"""
        log_file = self.output_dir / "extraction.log"
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(levelname)s - %(message)s',
            handlers=[
                logging.FileHandler(log_file),
                logging.StreamHandler()
            ]
        )
        self.logger = logging.getLogger(__name__)
    
    def _init_ocr_models(self):
        """Initialize OCR models if available"""
        if EASYOCR_AVAILABLE:
            try:
                self.ocr_reader = easyocr.Reader(['en'])
                self.logger.info("EasyOCR initialized successfully")
            except Exception as e:
                self.logger.warning(f"Failed to initialize EasyOCR: {e}")
        
        if PIX2TEX_AVAILABLE:
            try:
                self.latex_ocr = LatexOCR()
                self.logger.info("Pix2Tex LaTeX OCR initialized successfully")
            except Exception as e:
                self.logger.warning(f"Failed to initialize Pix2Tex: {e}")
    
    def _emu_to_pixels(self, emu_value: int) -> float:
        """Convert EMU (English Metric Units) to pixels"""
        return emu_value / 9525.0  # 1 inch = 914400 EMU, 1 inch = 96 pixels
    
    def _extract_layout_info(self, shape: BaseShape) -> LayoutInfo:
        """Extract layout information from a shape"""
        try:
            left = self._emu_to_pixels(shape.left)
            top = self._emu_to_pixels(shape.top)
            width = self._emu_to_pixels(shape.width)
            height = self._emu_to_pixels(shape.height)
            
            # Get rotation if available
            rotation = 0.0
            if hasattr(shape, 'rotation'):
                rotation = shape.rotation
            
            return LayoutInfo(
                left=left,
                top=top,
                width=width,
                height=height,
                rotation=rotation
            )
        except Exception as e:
            self.logger.warning(f"Failed to extract layout info: {e}")
            return LayoutInfo(0, 0, 100, 50)
    
    def _extract_text_formatting(self, shape: BaseShape) -> Optional[TextFormatting]:
        """Extract text formatting information"""
        try:
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                return None
            
            text_frame = shape.text_frame
            if not text_frame.paragraphs:
                return None
            
            # Get formatting from first paragraph/run
            paragraph = text_frame.paragraphs[0]
            if not paragraph.runs:
                return TextFormatting()
            
            run = paragraph.runs[0]
            font = run.font
            
            formatting = TextFormatting(
                font_name=font.name or "Arial",
                font_size=float(font.size.pt) if font.size else 12.0,
                bold=font.bold or False,
                italic=font.italic or False,
                underline=font.underline or False,
                alignment=str(paragraph.alignment) if paragraph.alignment else "left"
            )
            
            # Extract color if available
            if font.color and font.color.rgb:
                rgb = font.color.rgb
                formatting.color = f"#{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"
            
            return formatting
        except Exception as e:
            self.logger.warning(f"Failed to extract text formatting: {e}")
            return TextFormatting()
    
    def _save_image(self, image_blob: bytes, slide_num: int, shape_id: str, 
                   image_format: str = "png") -> str:
        """Save image to file and return relative path"""
        try:
            # Create unique filename
            image_hash = hashlib.md5(image_blob).hexdigest()[:8]
            filename = f"slide_{slide_num:03d}_{shape_id}_{image_hash}.{image_format}"
            image_path = self.images_dir / filename
            
            # Save image
            with open(image_path, 'wb') as f:
                f.write(image_blob)
            
            # Return relative path for markdown
            return f"images/{filename}"
        except Exception as e:
            self.logger.error(f"Failed to save image: {e}")
            return ""
    
    def _perform_ocr(self, image_path: str) -> Tuple[str, List[str]]:
        """Perform OCR on image and detect equations"""
        ocr_text = ""
        equations = []
        
        try:
            if not os.path.exists(image_path):
                return ocr_text, equations
            
            # Read image
            image = cv2.imread(image_path)  # type: ignore
            if image is None:
                return ocr_text, equations
            
            # Standard OCR with EasyOCR
            if self.ocr_reader:
                try:
                    results = self.ocr_reader.readtext(image)
                    ocr_text = " ".join([result[1] for result in results])
                except Exception as e:
                    self.logger.warning(f"EasyOCR failed: {e}")
            
            # Equation detection and LaTeX conversion
            if self.latex_ocr:
                try:
                    # Convert to PIL Image for pix2tex
                    pil_image = Image.fromarray(cv2.cvtColor(image, cv2.COLOR_BGR2RGB))  # type: ignore
                    latex_result = self.latex_ocr(pil_image)
                    
                    # Simple heuristic to detect if result looks like an equation
                    if self._is_likely_equation(latex_result):
                        equations.append(latex_result)
                except Exception as e:
                    self.logger.warning(f"LaTeX OCR failed: {e}")
            
        except Exception as e:
            self.logger.error(f"OCR processing failed: {e}")
        
        return ocr_text, equations
    
    def _is_likely_equation(self, latex_text: str) -> bool:
        """Heuristic to determine if LaTeX text represents an equation"""
        equation_indicators = [
            r'\\frac', r'\\sum', r'\\int', r'\\sqrt', r'\\alpha', r'\\beta',
            r'\\gamma', r'\\delta', r'\\theta', r'\\lambda', r'\\mu', r'\\sigma',
            r'\\pi', r'\\phi', r'\\psi', r'\\omega', r'\\partial', r'\\nabla',
            r'\\infty', r'\\pm', r'\\times', r'\\div', r'\\leq', r'\\geq',
            r'\\neq', r'\\approx', r'\\equiv', r'\\propto', r'\\in', r'\\subset',
            r'\\cup', r'\\cap', r'\\lim', r'\\log', r'\\ln', r'\\sin', r'\\cos',
            r'\\tan', r'\\exp', r'^{', r'_{', r'\\left', r'\\right'
        ]
        
        return any(indicator in latex_text for indicator in equation_indicators)
    
    def _extract_shape_data(self, shape: BaseShape, slide_num: int) -> ShapeData:
        """Extract complete data from a shape"""
        shape_id = f"shape_{id(shape)}"
        layout = self._extract_layout_info(shape)
        
        # Determine shape type
        shape_type = "unknown"
        if hasattr(shape, 'shape_type'):
            shape_type = str(shape.shape_type).split('.')[-1].lower()
        
        shape_data = ShapeData(
            shape_id=shape_id,
            shape_type=shape_type,
            layout=layout
        )
        
        try:
            # Handle different shape types
            if isinstance(shape, Picture):
                # Extract image
                image_blob = shape.image.blob
                image_path = self._save_image(image_blob, slide_num, shape_id)
                shape_data.image_path = image_path
                
                # Perform OCR on image
                if image_path:
                    full_image_path = self.output_dir / image_path
                    ocr_text, equations = self._perform_ocr(str(full_image_path))
                    shape_data.ocr_text = ocr_text
                    shape_data.latex_equations = equations
            
            elif hasattr(shape, 'text') and shape.text.strip():
                # Text shape
                shape_data.content = shape.text.strip()
                shape_data.formatting = self._extract_text_formatting(shape)
            
            elif hasattr(shape, 'table') and shape.table is not None:
                # Table shape
                shape_data.content = self._extract_table_content(shape.table)
                shape_data.shape_type = "table"
            
            elif isinstance(shape, GroupShape):
                # Group shape - extract from grouped shapes
                group_content = []
                for grouped_shape in shape.shapes:
                    grouped_data = self._extract_shape_data(grouped_shape, slide_num)
                    if grouped_data.content or grouped_data.image_path:
                        group_content.append(grouped_data.content or f"![Image]({grouped_data.image_path})")
                shape_data.content = "\n".join(group_content)
                shape_data.shape_type = "group"
        
        except Exception as e:
            self.logger.warning(f"Failed to extract shape data: {e}")
        
        return shape_data
    
    def _extract_table_content(self, table: Table) -> str:
        """Extract content from a table and format as markdown table"""
        if not table.rows:
            return ""
        
        table_lines = []
        
        # Extract header row
        if table.rows:
            header_cells = []
            for cell in table.rows[0].cells:
                header_cells.append(cell.text.strip() or " ")
            table_lines.append("| " + " | ".join(header_cells) + " |")
            table_lines.append("|" + "---|" * len(header_cells))
            
            # Extract data rows
            for row in table.rows[1:]:
                data_cells = []
                for cell in row.cells:
                    data_cells.append(cell.text.strip() or " ")
                table_lines.append("| " + " | ".join(data_cells) + " |")
        
        return "\n".join(table_lines)
    
    def extract_slide_content(self, slide, slide_num: int) -> Dict[str, Any]:
        """Extract complete slide content with layout preservation"""
        self.slide_counter += 1
        
        slide_content = {
            'slide_number': slide_num,
            'title': '',
            'shapes': [],
            'layout_data': {},
            'notes': '',
            'slide_size': {'width': 0, 'height': 0}
        }
        
        # Get slide dimensions
        if hasattr(slide, 'slide_layout') and hasattr(slide.slide_layout, 'slide_master'):
            master = slide.slide_layout.slide_master
            if hasattr(master, 'slide_width') and hasattr(master, 'slide_height'):
                slide_content['slide_size'] = {
                    'width': self._emu_to_pixels(master.slide_width),
                    'height': self._emu_to_pixels(master.slide_height)
                }
        
        # Extract all shapes with layout information
        title_found = False
        shapes_data = []
        
        for i, shape in enumerate(slide.shapes):
            try:
                shape_data = self._extract_shape_data(shape, slide_num)
                shape_data.layout.z_order = i  # Preserve stacking order
                
                # Try to identify title (first text shape at top of slide)
                if (not title_found and shape_data.content and 
                    shape_data.shape_type in ['text', 'textbox'] and
                    shape_data.layout.top < 100 and len(shape_data.content) < 100):
                    slide_content['title'] = shape_data.content
                    title_found = True
                
                shapes_data.append(shape_data)
                
            except Exception as e:
                self.logger.warning(f"Failed to process shape {i}: {e}")
                continue
        
        slide_content['shapes'] = shapes_data
        
        # Extract slide notes
        try:
            if hasattr(slide, 'notes_slide') and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_content['notes'] = notes_text
        except Exception as e:
            self.logger.warning(f"Failed to extract notes: {e}")
        
        return slide_content
    
    def extract_presentation(self, pptx_path: str) -> Dict[str, Any]:
        """Extract complete presentation with layout preservation"""
        try:
            prs = Presentation(pptx_path)
            self.presentation_counter += 1
            
            presentation_data = {
                'filename': Path(pptx_path).name,
                'title': Path(pptx_path).stem,
                'slides': [],
                'slide_count': len(prs.slides),
                'presentation_id': self.presentation_counter,
                'metadata': {
                    'slide_width': self._emu_to_pixels(prs.slide_width),
                    'slide_height': self._emu_to_pixels(prs.slide_height)
                }
            }
            
            # Reset slide counter for each presentation
            self.slide_counter = 0
            
            # Create progress bar if tqdm is available
            if TQDM_AVAILABLE:
                slide_iterator = tqdm(enumerate(prs.slides, 1), 
                                    total=len(prs.slides),
                                    desc=f"Processing {Path(pptx_path).name}",
                                    unit="slide")
            else:
                slide_iterator = enumerate(prs.slides, 1)
            
            for i, slide in slide_iterator:
                slide_content = self.extract_slide_content(slide, i)
                presentation_data['slides'].append(slide_content)
                if not TQDM_AVAILABLE:
                    self.logger.info(f"Processed slide {i}/{len(prs.slides)}")
            
            return presentation_data
            
        except Exception as e:
            self.logger.error(f"Error processing {pptx_path}: {str(e)}")
            return None


class PresentationGenerator:
    """Generate presentations from extracted markdown and JSON data"""
    
    def __init__(self, output_dir: str = "extracted_presentations"):
        self.output_dir = Path(output_dir)
        self.logger = logging.getLogger(__name__)
    
    def generate_clean_markdown(self, json_file: str, output_file: str = None, image_folder: str = "images") -> str:
        """Generate clean markdown with LaTeX equations and image references"""
        try:
            with open(json_file, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            if not output_file:
                json_path = Path(json_file)
                output_file = json_path.parent / f"{json_path.stem}_clean.md"
            
            clean_content = []
            
            # Handle both single presentation and list of presentations
            presentations = data if isinstance(data, list) else [data]
            
            for presentation in presentations:
                clean_content.append(f"# {presentation['title']}\n")
                
                for slide_idx, slide in enumerate(presentation['slides'], 1):
                    clean_content.append(f"## Slide {slide_idx}\n")
                    
                    # Process shapes - they are stored as string representations
                    for shape_str in slide['shapes']:
                        # Parse the string representation to extract useful information
                        if 'content=' in shape_str:
                            # Extract content between quotes after 'content='
                            import re
                            content_match = re.search(r"content='([^']*)'", shape_str)
                            if content_match:
                                content = content_match.group(1).strip()
                                if content and content != '':
                                    # Unescape common characters
                                    content = content.replace('\\n', '\n').replace('\\xa0', ' ')
                                    clean_content.append(f"{content}\n")
                        
                        # Extract image path
                        if 'image_path=' in shape_str:
                            image_match = re.search(r"image_path='([^']*)'", shape_str)
                            if image_match:
                                image_path = image_match.group(1)
                                if image_path and image_path != '':
                                    # Use the correct path to the images directory
                                    img_path = f"{image_folder}/{image_path}"
                                    clean_content.append(f"![Image]({img_path})\n")
                        
                        # Extract OCR text
                        if 'ocr_text=' in shape_str:
                            ocr_match = re.search(r"ocr_text='([^']*)'", shape_str)
                            if ocr_match:
                                ocr_text = ocr_match.group(1)
                                if ocr_text and ocr_text != '':
                                    clean_content.append(f"*OCR Text: {ocr_text}*\n")
                        
                        # Extract LaTeX equations
                        if 'latex_equations=' in shape_str and 'latex_equations=[]' not in shape_str:
                            # Extract LaTeX equations from the string representation
                            latex_match = re.search(r"latex_equations=\[(.*?)\]", shape_str, re.DOTALL)
                            if latex_match:
                                equations_str = latex_match.group(1)
                                # Split by quotes and clean up
                                equations = []
                                for eq in re.findall(r"'([^']*)'", equations_str):
                                    if eq.strip():
                                        # Clean up the LaTeX equation
                                        cleaned_eq = eq.replace('\\\\', '\\').strip()
                                        if cleaned_eq:
                                            equations.append(cleaned_eq)
                                
                                if equations:
                                    clean_content.append("**LaTeX Equations:**\n")
                                    for i, eq in enumerate(equations, 1):
                                        # Format as display math for better PDF rendering
                                        clean_content.append(f"$$\n{eq}\n$$\n")
                                    clean_content.append("\n")
                    
                    clean_content.append("\n---\n\n")
            
            # Write clean markdown
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write('\n'.join(clean_content))
            
            self.logger.info(f"Generated clean markdown: {output_file}")
            return str(output_file)
            
        except Exception as e:
            self.logger.error(f"Error generating clean markdown: {e}")
            return None
    
    def generate_pdf_presentation(self, markdown_file: str, json_file: str, output_file: str = None) -> str:
        """Generate PDF presentation from markdown and JSON data"""
        try:
            if not output_file:
                md_path = Path(markdown_file)
                output_file = md_path.parent / f"{md_path.stem}_presentation.pdf"
            
            # Try to use pandoc to convert markdown to PDF
            import subprocess
            
            try:
                # Check if pandoc is available
                subprocess.run(['pandoc', '--version'], capture_output=True, check=True)
                
                # First, create a cleaned version of the markdown for PDF generation
                cleaned_md_file = self._create_pdf_compatible_markdown(markdown_file)
                
                # Generate PDF using pandoc with safer options
                cmd = [
                    'pandoc', 
                    cleaned_md_file,
                    '-o', str(output_file),
                    '--pdf-engine=pdflatex',
                    '-V', 'geometry:margin=1in',
                    '--from=markdown',
                    '--to=pdf'
                ]
                
                result = subprocess.run(cmd, capture_output=True, text=True)
                
                # Clean up temporary file
                if cleaned_md_file != markdown_file:
                    Path(cleaned_md_file).unlink(missing_ok=True)
                
                if result.returncode == 0:
                    self.logger.info(f"PDF presentation generated: {output_file}")
                    print(f"✅ PDF presentation generated: {output_file}")
                    return str(output_file)
                else:
                    self.logger.error(f"Pandoc error: {result.stderr}")
                    print(f"❌ Pandoc error: {result.stderr}")
                    return None
                    
            except (subprocess.CalledProcessError, FileNotFoundError):
                self.logger.warning("Pandoc not found. Please install pandoc to generate PDF presentations.")
                print("❌ Pandoc not found. Please install pandoc to generate PDF presentations.")
                print("   Install with: brew install pandoc (macOS) or apt-get install pandoc (Ubuntu)")
                return None
            
        except Exception as e:
            self.logger.error(f"Error generating PDF: {e}")
            return None
    
    def _create_pdf_compatible_markdown(self, markdown_file: str) -> str:
        """Create a PDF-compatible version of the markdown file"""
        try:
            with open(markdown_file, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Clean up problematic content for PDF generation
            lines = content.split('\n')
            cleaned_lines = []
            skip_latex_block = False
            
            for line in lines:
                # Skip LaTeX equation blocks entirely
                if line.strip() == '**LaTeX Equations:**':
                    skip_latex_block = True
                    continue
                elif line.strip().startswith('$$'):
                    if skip_latex_block:
                        continue
                    else:
                        skip_latex_block = not skip_latex_block
                        continue
                elif skip_latex_block and line.strip() == '':
                    skip_latex_block = False
                    continue
                elif skip_latex_block:
                    continue
                
                # Skip LaTeX equation placeholders
                if '*[LaTeX equations detected]*' in line:
                    continue
                
                # Clean up image paths for PDF
                if line.startswith('![Image]('):
                    # Skip images that might not be found
                    continue
                
                # Clean up OCR text formatting and handle Unicode characters
                if line.startswith('*OCR Text:'):
                    # Keep OCR text but format it better and replace problematic Unicode
                    ocr_text = line.replace('*OCR Text:', '').replace('*', '').strip()
                    if ocr_text:
                        # Replace common Unicode characters that cause LaTeX issues
                        ocr_text = ocr_text.replace('≥', '>=').replace('≤', '<=').replace('≠', '!=')
                        ocr_text = ocr_text.replace('→', '->').replace('←', '<-').replace('↔', '<->')
                        ocr_text = ocr_text.replace('∞', 'infinity').replace('±', '+/-')
                        ocr_text = ocr_text.replace('×', 'x').replace('÷', '/')
                        ocr_text = ocr_text.replace('°', ' degrees').replace('²', '^2').replace('³', '^3')
                        ocr_text = ocr_text.replace('∈', ' in ').replace('∉', ' not in ')
                        ocr_text = ocr_text.replace('∪', ' union ').replace('∩', ' intersection ')
                        ocr_text = ocr_text.replace('⊂', ' subset ').replace('⊃', ' superset ')
                        ocr_text = ocr_text.replace('∀', ' for all ').replace('∃', ' exists ')
                        ocr_text = ocr_text.replace('∑', ' sum ').replace('∏', ' product ')
                        ocr_text = ocr_text.replace('∫', ' integral ').replace('∂', ' partial ')
                        cleaned_lines.append(f"**OCR Text:** {ocr_text}")
                else:
                    # Handle Unicode characters in regular text as well
                    cleaned_line = line.replace('≥', '>=').replace('≤', '<=').replace('≠', '!=')
                    cleaned_line = cleaned_line.replace('→', '->').replace('←', '<-').replace('↔', '<->')
                    cleaned_line = cleaned_line.replace('∞', 'infinity').replace('±', '+/-')
                    cleaned_line = cleaned_line.replace('×', 'x').replace('÷', '/')
                    cleaned_line = cleaned_line.replace('°', ' degrees').replace('²', '^2').replace('³', '^3')
                    cleaned_line = cleaned_line.replace('∈', ' in ').replace('∉', ' not in ')
                    cleaned_line = cleaned_line.replace('∪', ' union ').replace('∩', ' intersection ')
                    cleaned_line = cleaned_line.replace('⊂', ' subset ').replace('⊃', ' superset ')
                    cleaned_line = cleaned_line.replace('∀', ' for all ').replace('∃', ' exists ')
                    cleaned_line = cleaned_line.replace('∑', ' sum ').replace('∏', ' product ')
                    cleaned_line = cleaned_line.replace('∫', ' integral ').replace('∂', ' partial ')
                    cleaned_lines.append(cleaned_line)
            
            # Create temporary cleaned file
            temp_file = Path(markdown_file).parent / f"temp_cleaned_{Path(markdown_file).name}"
            with open(temp_file, 'w', encoding='utf-8') as f:
                f.write('\n'.join(cleaned_lines))
            
            return str(temp_file)
            
        except Exception as e:
            self.logger.error(f"Error creating PDF-compatible markdown: {e}")
            return markdown_file


def find_pptx_files(folder_path: str) -> List[str]:
    """Find all PPTX files in the specified folder"""
    folder = Path(folder_path)
    if not folder.exists() or not folder.is_dir():
        print(f"Error: Folder '{folder_path}' does not exist or is not a directory")
        return []
    
    pptx_files = list(folder.glob("*.pptx"))
    pptx_files.extend(list(folder.glob("*.PPTX")))
    
    # Sort files naturally
    def natural_sort_key(filename):
        return [int(text) if text.isdigit() else text.lower() 
                for text in re.split(r'(\d+)', str(filename))]
    
    pptx_files.sort(key=natural_sort_key)
    return [str(f) for f in pptx_files]


def format_shape_as_markdown(shape_data: ShapeData, include_layout: bool = True) -> str:
    """Format a shape as presentation-ready markdown"""
    markdown_lines = []
    
    # Add layout information as HTML comment for presentation tools
    if include_layout:
        layout_info = {
            'position': {
                'left': shape_data.layout.left,
                'top': shape_data.layout.top,
                'width': shape_data.layout.width,
                'height': shape_data.layout.height,
                'rotation': shape_data.layout.rotation,
                'z_order': shape_data.layout.z_order
            },
            'type': shape_data.shape_type
        }
        markdown_lines.append(f"<!-- LAYOUT: {json.dumps(layout_info)} -->")
    
    # Handle different content types
    if shape_data.image_path:
        # Image with optional OCR text and equations
        markdown_lines.append(f"![Image]({shape_data.image_path})")
        
        if shape_data.ocr_text:
            markdown_lines.append(f"<!-- OCR: {shape_data.ocr_text} -->")
        
        if shape_data.latex_equations:
            for eq in shape_data.latex_equations:
                markdown_lines.append(f"$$\n{eq}\n$$")
    
    elif shape_data.content:
        # Text content with formatting
        content = shape_data.content
        
        # Apply formatting if available
        if shape_data.formatting:
            fmt = shape_data.formatting
            if fmt.bold:
                content = f"**{content}**"
            if fmt.italic:
                content = f"*{content}*"
            if fmt.font_size > 18:  # Large text, likely a header
                content = f"## {content}"
            elif fmt.font_size > 14:  # Medium text, subheader
                content = f"### {content}"
        
        # Handle tables
        if shape_data.shape_type == "table":
            markdown_lines.append(content)
        else:
            # Handle bullet points and regular text
            lines = content.split('\n')
            for line in lines:
                line = line.strip()
                if line:
                    if line.startswith('•') or line.startswith('-'):
                        markdown_lines.append(f"- {line[1:].strip()}")
                    else:
                        markdown_lines.append(line)
    
    return "\n".join(markdown_lines)


def format_slide_as_presentation_markdown(slide_data: Dict[str, Any], 
                                        presentation_title: str,
                                        include_layout: bool = True) -> str:
    """Format a slide as presentation-ready markdown"""
    markdown_lines = []
    
    # Slide separator for presentation tools
    markdown_lines.append("---")
    markdown_lines.append("")
    
    # Slide metadata
    if include_layout:
        slide_meta = {
            'slide_number': slide_data['slide_number'],
            'slide_size': slide_data['slide_size'],
            'shape_count': len(slide_data['shapes'])
        }
        markdown_lines.append(f"<!-- SLIDE_META: {json.dumps(slide_meta)} -->")
        markdown_lines.append("")
    
    # Slide title
    if slide_data['title']:
        markdown_lines.append(f"# {slide_data['title']}")
    else:
        markdown_lines.append(f"# Slide {slide_data['slide_number']}")
    
    markdown_lines.append("")
    
    # Sort shapes by z-order and position for logical reading order
    shapes = sorted(slide_data['shapes'], 
                   key=lambda s: (s.layout.z_order, s.layout.top, s.layout.left))
    
    # Process each shape
    for shape in shapes:
        if shape.content or shape.image_path:
            shape_markdown = format_shape_as_markdown(shape, include_layout)
            if shape_markdown.strip():
                markdown_lines.append(shape_markdown)
                markdown_lines.append("")
    
    # Add speaker notes
    if slide_data['notes']:
        markdown_lines.append("<!-- SPEAKER_NOTES")
        markdown_lines.append(slide_data['notes'])
        markdown_lines.append("-->")
        markdown_lines.append("")
    
    return "\n".join(markdown_lines)


def create_presentation_markdown(presentations: List[Dict[str, Any]], 
                               output_file: str,
                               include_layout: bool = True):
    """Create presentation-ready markdown file"""
    markdown_content = []
    
    # Document header with metadata
    markdown_content.append("---")
    markdown_content.append("title: Combined Presentations")
    markdown_content.append("theme: default")
    markdown_content.append("transition: slide")
    markdown_content.append("highlightTheme: github")
    markdown_content.append("---")
    markdown_content.append("")
    
    # Title slide
    markdown_content.append("# Combined Presentation Content")
    markdown_content.append("")
    markdown_content.append(f"*Generated from {len(presentations)} PowerPoint presentations*")
    markdown_content.append("")
    
    # Process each presentation
    for pres in presentations:
        # Presentation separator
        markdown_content.append("---")
        markdown_content.append("")
        markdown_content.append(f"# {pres['title']}")
        markdown_content.append("")
        markdown_content.append(f"*Source: {pres['filename']}*")
        markdown_content.append("")
        
        # Process each slide
        for slide in pres['slides']:
            slide_markdown = format_slide_as_presentation_markdown(
                slide, pres['title'], include_layout
            )
            markdown_content.append(slide_markdown)
    
    # Write to file
    try:
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write("\n".join(markdown_content))
        print(f"✅ Presentation-ready markdown created: {output_file}")
        
        # Also create a JSON file with raw data
        json_file = output_file.replace('.md', '_data.json')
        with open(json_file, 'w', encoding='utf-8') as f:
            json.dump(presentations, f, indent=2, default=str)
        print(f"✅ Raw data JSON created: {json_file}")
        
    except Exception as e:
        print(f"❌ Error writing files: {str(e)}")


def action_extract(args):
    """Extract PPTX files to markdown and JSON"""
    print("🚀 Enhanced PPTX Extractor Starting...")
    print(f"🔍 Input folder: {args.input_folder}")
    print(f"📝 Output file: {args.output_file}")
    print(f"📁 Output directory: {args.output_dir}")
    
    # Check OCR availability
    if not args.no_ocr:
        if EASYOCR_AVAILABLE:
            print("✅ EasyOCR available")
        else:
            print("⚠️  EasyOCR not available - install with: pip install easyocr")
        
        if PIX2TEX_AVAILABLE:
            print("✅ Pix2Tex LaTeX OCR available")
        else:
            print("⚠️  Pix2Tex not available - install with: pip install pix2tex")
    
    print("")
    
    # Find PPTX files
    pptx_files = find_pptx_files(args.input_folder)
    
    if not pptx_files:
        print("❌ No PPTX files found")
        sys.exit(1)
    
    print(f"📁 Found {len(pptx_files)} PPTX files:")
    for i, file in enumerate(pptx_files, 1):
        print(f"  {i}. {Path(file).name}")
    print("")
    
    # Initialize extractor
    extractor = EnhancedPPTXExtractor(args.output_dir)
    presentations = []
    
    # Process each presentation
    for pptx_file in pptx_files:
        print(f"📖 Processing: {Path(pptx_file).name}")
        presentation_data = extractor.extract_presentation(pptx_file)
        
        if presentation_data:
            presentations.append(presentation_data)
            print(f"   ✅ Extracted {presentation_data['slide_count']} slides")
            
            # Count images and equations
            image_count = sum(1 for slide in presentation_data['slides'] 
                            for shape in slide['shapes'] if shape.image_path)
            equation_count = sum(len(shape.latex_equations) 
                               for slide in presentation_data['slides'] 
                               for shape in slide['shapes'])
            
            if image_count > 0:
                print(f"   🖼️  Extracted {image_count} images")
            if equation_count > 0:
                print(f"   📐 Detected {equation_count} equations")
                
            # Generate individual markdown file for each PPTX
            pptx_path = Path(pptx_file)
            if args.output_file:
                # If output file is specified and there's only one PPTX, use it
                if len(pptx_files) == 1:
                    output_filename = extractor.output_dir / args.output_file
                else:
                    # Multiple files: use PPTX name with specified extension
                    output_filename = extractor.output_dir / f"{pptx_path.stem}.md"
            else:
                # Default: use PPTX filename with .md extension in output directory
                output_filename = extractor.output_dir / f"{pptx_path.stem}.md"
            
            print(f"   📝 Creating markdown: {output_filename}")
            
            # Save JSON data alongside markdown
            json_filename = output_filename.with_suffix('.json')
            with open(json_filename, 'w', encoding='utf-8') as f:
                json.dump(presentation_data, f, indent=2, ensure_ascii=False, default=str)
            print(f"   💾 Saved data: {json_filename}")
            
            # Create individual markdown for this presentation
            create_presentation_markdown(
                [presentation_data], 
                str(output_filename),
                include_layout=not args.no_layout
            )
            
        else:
            print(f"   ❌ Failed to process")
    
    if not presentations:
        print("❌ No presentations processed successfully")
        sys.exit(1)

    print("")
    
    # Summary
    total_slides = sum(pres['slide_count'] for pres in presentations)
    total_images = sum(1 for pres in presentations 
                      for slide in pres['slides'] 
                      for shape in slide['shapes'] if shape.image_path)
    total_equations = sum(len(shape.latex_equations) 
                         for pres in presentations 
                         for slide in pres['slides'] 
                         for shape in slide['shapes'])
    
    print("📊 Summary:")
    print(f"   - Presentations: {len(presentations)}")
    print(f"   - Total slides: {total_slides}")
    print(f"   - Images extracted: {total_images}")
    print(f"   - Equations detected: {total_equations}")
    print(f"   - Output directory: {args.output_dir}")
    print("")
    print("🎉 Enhanced extraction completed!")


def action_generate_pdf(args):
    """Generate PDF presentation from markdown and JSON files"""
    print("📄 PDF Generation Starting...")
    
    if not args.json_file:
        print("❌ JSON file is required for PDF generation")
        sys.exit(1)
    
    json_path = Path(args.json_file)
    if not json_path.exists():
        print(f"❌ JSON file not found: {args.json_file}")
        sys.exit(1)
    
    # Find corresponding markdown file if not specified
    markdown_file = args.markdown_file
    if not markdown_file:
        markdown_file = json_path.with_suffix('.md')
        if not markdown_file.exists():
            print(f"❌ Corresponding markdown file not found: {markdown_file}")
            sys.exit(1)
    
    generator = PresentationGenerator(args.output_dir)
    pdf_file = generator.generate_pdf_presentation(
        str(markdown_file), 
        str(json_path), 
        args.output_file
    )
    
    if pdf_file:
        print(f"✅ PDF generation completed: {pdf_file}")
    else:
        print("❌ PDF generation failed")
        sys.exit(1)


def action_clean_markdown(args):
    """Generate clean markdown with LaTeX equations and image references"""
    print("🧹 Clean Markdown Generation Starting...")
    
    if not args.json_file:
        print("❌ JSON file is required for clean markdown generation")
        sys.exit(1)
    
    json_path = Path(args.json_file)
    if not json_path.exists():
        print(f"❌ JSON file not found: {args.json_file}")
        sys.exit(1)
    
    generator = PresentationGenerator(args.output_dir)
    clean_md_file = generator.generate_clean_markdown(
        str(json_path), 
        args.output_file
    )
    
    if clean_md_file:
        print(f"✅ Clean markdown generated: {clean_md_file}")
    else:
        print("❌ Clean markdown generation failed")
        sys.exit(1)


def main():
    """Main function with action-based commands"""
    parser = argparse.ArgumentParser(
        description="Enhanced PPTX to Markdown converter with layout preservation and OCR",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""Examples:
  # Extract PPTX to markdown (default action)
  python enhanced_pptx_extractor.py extract /path/to/pptx/folder
  python enhanced_pptx_extractor.py extract /path/to/pptx/folder --output-file presentation.md
  
  # Generate PDF from extracted data
  python enhanced_pptx_extractor.py generate-pdf --json-file output/presentation.json
  
  # Generate clean markdown
  python enhanced_pptx_extractor.py clean-markdown --json-file output/presentation.json"""
    )
    
    # Add subcommands
    subparsers = parser.add_subparsers(dest='action', help='Available actions')
    
    # Extract command (default)
    extract_parser = subparsers.add_parser('extract', help='Extract PPTX files to markdown and JSON')
    extract_parser.add_argument(
        'input_folder',
        nargs='?',
        default='data/pptx',
        help='Path to folder containing PPTX files (default: data/pptx)'
    )
    extract_parser.add_argument(
        '--output-file',
        help='Output markdown file (default: same as PPTX filename with .md extension)'
    )
    extract_parser.add_argument(
        '--output-dir',
        default='output/extracted_presentations',
        help='Output directory for images and data (default: output/extracted_presentations)'
    )
    extract_parser.add_argument(
        '--no-layout',
        action='store_true',
        help='Exclude layout information from output'
    )
    extract_parser.add_argument(
        '--no-ocr',
        action='store_true',
        help='Disable OCR processing'
    )
    
    # Generate PDF command
    pdf_parser = subparsers.add_parser('generate-pdf', help='Generate PDF presentation from markdown and JSON')
    pdf_parser.add_argument(
        '--json-file',
        required=True,
        help='JSON file containing presentation data'
    )
    pdf_parser.add_argument(
        '--markdown-file',
        help='Markdown file (default: inferred from JSON filename)'
    )
    pdf_parser.add_argument(
        '--output-file',
        help='Output PDF file (default: inferred from JSON filename)'
    )
    pdf_parser.add_argument(
        '--output-dir',
        default='output/extracted_presentations',
        help='Output directory (default: output/extracted_presentations)'
    )
    
    # Clean markdown command
    clean_parser = subparsers.add_parser('clean-markdown', help='Generate clean markdown with LaTeX and image references')
    clean_parser.add_argument(
        '--json-file',
        required=True,
        help='JSON file containing presentation data'
    )
    clean_parser.add_argument(
        '--output-file',
        help='Output clean markdown file (default: inferred from JSON filename)'
    )
    clean_parser.add_argument(
        '--output-dir',
        default='output/extracted_presentations',
        help='Output directory (default: output/extracted_presentations)'
    )
    
    args = parser.parse_args()
    
    # Default action is extract if no action specified
    if not args.action:
        # Handle legacy usage (no subcommand)
        args.action = 'extract'
        # Parse remaining args as extract arguments
        if len(sys.argv) > 1 and not sys.argv[1].startswith('-'):
            args.input_folder = sys.argv[1]
        else:
            args.input_folder = 'data/pptx'
        args.output_file = sys.argv[2] if len(sys.argv) > 2 and not sys.argv[2].startswith('-') else None
        args.output_dir = 'output/extracted_presentations'
        args.no_layout = '--no-layout' in sys.argv
        args.no_ocr = '--no-ocr' in sys.argv
    
    # Execute the appropriate action
    if args.action == 'extract':
        action_extract(args)
    elif args.action == 'generate-pdf':
        action_generate_pdf(args)
    elif args.action == 'clean-markdown':
        action_clean_markdown(args)
    else:
        parser.print_help()
        sys.exit(1)


if __name__ == "__main__":
    main()